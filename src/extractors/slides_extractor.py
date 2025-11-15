from pathlib import Path
from typing import List, Dict
from pptx import Presentation

class SlidesExtractor:
    """Extract and structure content from PowerPoint .pptx files"""
    
    def __init__(self):
        self.supported_format = '.pptx'
    
    def extract_from_file(self, file_path: Path) -> List[Dict]:
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if file_path.suffix != self.supported_format:
            raise ValueError(f"Unsupported format: {file_path.suffix}")
        
        prs = Presentation(file_path)
        slides_data = []
        
        for idx, slide in enumerate(prs.slides, start=1):
            slide_data = self._extract_slide_content(slide, idx)
            slides_data.append(slide_data)
        
        return slides_data
    
    def _extract_slide_content(self, slide, slide_number: int) -> Dict:
        
        # Content output for each slide
        content = {
            'slide_number': slide_number,
            'title': '',
            'text_content': [],
            'full_text': ''
        }
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                
                if not content['title'] and shape.shape_type == 1:  # Title shape
                    content['title'] = text
                else:
                    content['text_content'].append(text)
        
        full_text_parts = []
        if content['title']:
            full_text_parts.append(f"Title: {content['title']}")
        if content['text_content']:
            full_text_parts.append("Content: " + " ".join(content['text_content']))
        
        content['full_text'] = "\n".join(full_text_parts)
        
        return content

if __name__ == "__main__":
    extractor = SlidesExtractor()
    content = extractor.extract_from_file(Path("data/inputs/Intro to Java.pptx"))
    for slide in content:
        print(f"Slide {slide['slide_number']}:\n{slide['full_text']}\n")
